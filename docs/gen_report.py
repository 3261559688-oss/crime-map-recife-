#!/usr/bin/env python3
"""生成项目汇报：Word 文档 + 一页 PPT"""
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PI, Pt as PP, Emu
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUT_DIR = Path(__file__).parent
URL = "https://crime-map-recife.vercel.app"

# =================================================================
# 1) 生成 Word 文档
# =================================================================
def make_docx():
    doc = Document()

    # 全局样式
    style = doc.styles['Normal']
    style.font.name = 'PingFang SC'
    style.font.size = Pt(10.5)

    def H1(text):
        p = doc.add_heading(text, level=1)
        for r in p.runs:
            r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        return p

    def H2(text):
        p = doc.add_heading(text, level=2)
        for r in p.runs:
            r.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
        return p

    def P(text, bold=False, color=None, size=None):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.font.size = size or Pt(10.5)
        r.bold = bold
        if color: r.font.color.rgb = color
        return p

    # === 封面标题 ===
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run('巴西全国犯罪事件实时地图')
    r.font.size = Pt(24); r.bold = True
    r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = sub.add_run('MVP 项目汇报 · Crime Map Brasil')
    r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    link = doc.add_paragraph()
    link.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = link.add_run(f'在线访问：{URL}')
    r.font.size = Pt(11); r.italic = True
    r.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)

    doc.add_paragraph()

    # === 一、项目背景 ===
    H1('一、项目背景与目标')
    P('业务洞察：当前各类新闻 App 在巴西市场存在「治安信息分散、用户感知滞后」的痛点。'
      '我们做了一个 MVP 验证：把巴西全国的治安/犯罪类新闻实时聚合到一张地图上，'
      '让用户一眼看到自己周边发生了什么。')
    P('MVP 目标：用最低成本（$0/月）跑通"新闻抓取 → 智能分类 → 地图可视化 → 实时更新"全链路，'
      '验证产品价值。', bold=True)

    # === 二、竞品调研 ===
    H1('二、竞品调研')
    H2('🇺🇸 NewsBreak（核心对标）')
    table = doc.add_table(rows=1, cols=3)
    table.style = 'Light Grid Accent 1'
    hdr = table.rows[0].cells
    hdr[0].text = '维度'; hdr[1].text = 'NewsBreak'; hdr[2].text = '我们的认知'
    rows = [
        ('MAU', '4500 万+', '美国本土第二大新闻 App'),
        ('覆盖', '41,000+ 美国邮编', '极致下沉到社区'),
        ('更新量', '~100 万篇/天（41,000 篇/小时）', '大量是 AI 重写'),
        ('核心价值', '"Local news that matters to you"', '超本地化推送'),
        ('争议点', 'AI 假新闻、版权纠纷', '高风险'),
    ]
    for k, v, c in rows:
        cells = table.add_row().cells
        cells[0].text = k; cells[1].text = v; cells[2].text = c

    H2('我们提炼的判断')
    for line in [
        '1. "超本地化"是真需求 — 用户对 5km 内发生的事比国际新闻更关心',
        '2. AI 内容生成有版权/合规风险 — 我们 V1 不重写，只聚合并标注来源',
        '3. 巴西市场尚无对标产品 — 本地媒体分散，无人做全国级实时地图',
        '4. 治安话题是高粘性入口 — 枪击/抢劫天然是用户每天必看的"近邻信息"',
    ]:
        P(line)

    # === 三、技术方案 ===
    H1('三、技术方案')
    P('整体架构（全部 $0 成本）：', bold=True)
    P('   153 个 RSS 源 → GitHub Actions（每 30 分钟）→ Python 抓取/分类/去重')
    P('       → git push → Vercel 自动部署 → 全球 CDN → 用户访问')
    doc.add_paragraph()
    P('关键技术点：', bold=True)
    for line in [
        '• 数据源：153 个 RSS（G1 全 27 州 + G1 城市级 + R7 + UOL + Folha + Estadão + 各地小报 + BBC/CNN Brasil）',
        '• 抓取：Python ThreadPoolExecutor 20 线程并行，全量抓取 < 30 秒',
        '• 分类：60+ 关键词 + 11 种犯罪类型映射',
        '• 去重：URL 哈希 + 7 天滚动窗口',
        '• 地理：60+ 城市坐标库，URL/标题双重定位',
        '• 前端：原生 HTML + Leaflet 地图，无框架',
        '• 更新：GitHub Actions cron */30，自动 commit + Vercel 自动部署',
    ]:
        P(line)

    # === 四、做了什么 ===
    H1('四、做了什么（动作流水）')
    stages = [
        ('阶段 1：从 0 到 1', '调通 G1 RSS 源 → 数据 84 条'),
        ('阶段 2：扩量', '接入 27 个州 G1 → 数据 366 → 657 条'),
        ('阶段 3：实时化', 'GitHub Actions 每 30 分钟自动跑，端到端延迟 ≤ 30 分钟'),
        ('阶段 4：再扩量', 'RSS 扩到 118 个，引入 20 线程并行 → 数据破 1048 条'),
        ('阶段 5：精细化分类', '类型从 5 种扩到 11 种，关键词扩到 60+，修复 quadrilha junina 误判'),
        ('阶段 6：再扩量', 'RSS 扩到 153 个 → 数据破 2042 条'),
        ('阶段 7：UI/UX', '浅色主题、时间标签、双层筛选、抽屉式列表'),
        ('阶段 8：移动端适配', '响应式布局，修复假状态栏，全屏沉浸式'),
    ]
    for s, d in stages:
        P(s + '：' + d)

    # === 五、当前功能 ===
    H1('五、当前具备的功能')
    H2('🗺️ 地图能力')
    for line in [
        '• 巴西全境地图（CartoDB Light 浅色底）',
        '• 2042+ 个事件标记，11 种颜色 + emoji 区分',
        '• 点击 marker 弹窗：标题 / 城市 / 来源 / 时间 / 阅读原文直达',
    ]: P(line)
    H2('📊 数据能力')
    for line in [
        '• 153 个 RSS 源，覆盖 27 个州 / 60+ 城市',
        '• 11 种犯罪类型：杀人 / 抢劫 / 盗窃 / 性犯罪 / 毒品 / 绑架 / 家暴 / 警方 / 黑帮 / 诈骗 / 车辆',
        '• 时间窗：最近 7 天滚动，自动去重',
    ]: P(line)
    H2('⚡ 实时能力')
    for line in [
        '• 每 30 分钟自动抓取、自动部署',
        '• 端到端延迟 < 30 分钟（接近 RSS 物理极限）',
    ]: P(line)
    H2('🎨 体验能力')
    for line in [
        '• 浅色主题，移动端全屏沉浸',
        '• 城市/类型双层筛选，底部抽屉列表',
        '• 时间标签彩色编码（há 2h / ontem / há 3 dias）',
    ]: P(line)

    # === 六、核心指标 ===
    H1('六、关键数据指标')
    t = doc.add_table(rows=1, cols=2)
    t.style = 'Light Grid Accent 1'
    t.rows[0].cells[0].text = '指标'
    t.rows[0].cells[1].text = '数值'
    metrics = [
        ('数据规模', '2042 条犯罪事件'),
        ('媒体覆盖', '153 个 RSS / 58 个稳定源'),
        ('地理覆盖', '27 个州 / 60+ 城市（巴西全境）'),
        ('分类粒度', '11 种类型 + 60+ 关键词'),
        ('更新频率', '每 30 分钟自动'),
        ('端到端延迟', '< 30 分钟'),
        ('数据增长', '84 → 2042（+2330%）'),
        ('运营成本', '$0/月（GitHub Actions + Vercel 免费额度）'),
    ]
    for k, v in metrics:
        c = t.add_row().cells
        c[0].text = k; c[1].text = v

    # === 七、结论 ===
    H1('七、核心结论 & 下一步')
    H2('✅ 验证成果')
    for line in [
        '1. 技术可行 — $0 成本跑通全链路，2000+ 数据规模、30 分钟实时更新',
        '2. 数据可用 — 巴西公开媒体足够丰富，无需爬虫，合规',
        '3. 分类有效 — 基础关键词覆盖 78% 数据，剩 22% 进 outros 兜底',
    ]: P(line)
    H2('🚀 下一步建议')
    for line in [
        '🅰 高优：LLM 分类替代关键词（准确率 78% → 95%）',
        '🅰 高优：移动端推送（PWA + Push Notification）',
        '🅱 中优：接入 Twitter/X 警方官方账号（5 分钟级实时）',
        '🅱 中优：用户上报功能（UGC，建立数据壁垒）',
        '🅲 探索：AI 摘要 / 多语言（葡 → 英/中）',
    ]: P(line)
    H2('💰 商业潜力')
    for line in [
        '• 广告：本地商家按邮编投放（餐饮、健身、地产）',
        '• B2B：保险公司（车险定价）、地产中介（社区评估）数据接口',
        '• 订阅：高级用户「我家 1km 内的实时推送」',
    ]: P(line)

    # 落款
    doc.add_paragraph()
    end = doc.add_paragraph()
    end.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = end.add_run(f'演示链接：{URL}')
    r.font.size = Pt(11); r.bold = True
    r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

    out = OUT_DIR / '巴西犯罪地图_项目汇报.docx'
    doc.save(out)
    print(f'✅ Word: {out}')


# =================================================================
# 2) 生成一页 PPT
# =================================================================
def make_pptx():
    prs = Presentation()
    prs.slide_width = PI(13.333)
    prs.slide_height = PI(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 背景：深红 → 黑色 渐变（用矩形模拟）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRGB(0x0A, 0x0E, 0x1A)
    bg.line.fill.background()

    # 顶部红色色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PI(0.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRGB(0xDC, 0x26, 0x26)
    bar.line.fill.background()

    def add_text(left, top, width, height, text, size, bold=False, color=(0xFF,0xFF,0xFF), align='left'):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
        r = p.add_run(); r.text = text
        r.font.size = PP(size); r.font.bold = bold
        r.font.color.rgb = PRGB(*color)
        r.font.name = 'PingFang SC'
        return tb

    # === 标题区 ===
    add_text(PI(0.5), PI(0.35), PI(9), PI(0.6),
             '🗺️  巴西全国犯罪事件实时地图', 30, bold=True, color=(0xFF,0xFF,0xFF))
    add_text(PI(0.5), PI(0.95), PI(9), PI(0.4),
             'Crime Map Brasil  ·  MVP 项目汇报', 14, color=(0xFB,0xBF,0x24))

    add_text(PI(9.5), PI(0.4), PI(3.5), PI(0.4),
             '🌐 crime-map-recife.vercel.app', 11, color=(0x60,0xA5,0xFA), align='right')
    add_text(PI(9.5), PI(0.85), PI(3.5), PI(0.4),
             '运营成本 $0/月', 11, bold=True, color=(0x22,0xC5,0x5E), align='right')

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PI(0.5), PI(1.45), PI(12.3), Emu(20000))
    line.fill.solid(); line.fill.fore_color.rgb = PRGB(0xDC, 0x26, 0x26)
    line.line.fill.background()

    # === 三栏布局：竞品 / 我们 / 数据 ===
    col_top = PI(1.7)
    col_h = PI(2.6)
    col_w = PI(4.0)

    # 列 1：竞品
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PI(0.5), col_top, col_w, col_h)
    box1.fill.solid(); box1.fill.fore_color.rgb = PRGB(0x1F, 0x29, 0x37)
    box1.line.color.rgb = PRGB(0x37, 0x41, 0x51); box1.line.width = Emu(12700)

    add_text(PI(0.7), col_top + PI(0.15), col_w - PI(0.4), PI(0.4),
             '🇺🇸 竞品：NewsBreak', 16, bold=True, color=(0xFB,0xBF,0x24))

    tb = slide.shapes.add_textbox(PI(0.7), col_top + PI(0.7), col_w - PI(0.4), col_h - PI(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    items = [
        ('MAU', '4500 万 (美国第 2 大新闻 App)'),
        ('覆盖', '41,000+ 邮编 → 极致本地化'),
        ('更新', '100 万篇/天 (大量 AI 重写)'),
        ('风险', '版权纠纷 / AI 假新闻'),
        ('启示', '巴西本地化是空白市场'),
    ]
    for i, (k, v) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f'• {k}：'
        r.font.size = PP(11); r.font.bold = True
        r.font.color.rgb = PRGB(0xFB,0xBF,0x24); r.font.name = 'PingFang SC'
        r2 = p.add_run(); r2.text = v
        r2.font.size = PP(11); r2.font.color.rgb = PRGB(0xE5,0xE7,0xEB); r2.font.name = 'PingFang SC'
        p.space_after = PP(4)

    # 列 2：我们做了什么
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PI(4.65), col_top, col_w, col_h)
    box2.fill.solid(); box2.fill.fore_color.rgb = PRGB(0x7F, 0x1D, 0x1D)
    box2.line.color.rgb = PRGB(0xDC, 0x26, 0x26); box2.line.width = Emu(12700)

    add_text(PI(4.85), col_top + PI(0.15), col_w - PI(0.4), PI(0.4),
             '🚀 我们做了什么', 16, bold=True, color=(0xFF,0xFF,0xFF))

    tb = slide.shapes.add_textbox(PI(4.85), col_top + PI(0.7), col_w - PI(0.4), col_h - PI(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    items = [
        '聚合 153 个 RSS 源 (G1/UOL/Folha+地方报)',
        '20 线程并行抓取，30 秒抓完全量',
        '11 类犯罪 + 60 关键词智能分类',
        '7 天滚动窗口，URL 哈希去重',
        'GitHub Actions 每 30 分钟自动更新',
        'Vercel 自动部署，全球 CDN',
        '响应式前端，PC/手机沉浸体验',
    ]
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = '✓ ' + line
        r.font.size = PP(11); r.font.color.rgb = PRGB(0xFF,0xFF,0xFF); r.font.name = 'PingFang SC'
        p.space_after = PP(4)

    # 列 3：数据成果
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PI(8.8), col_top, col_w, col_h)
    box3.fill.solid(); box3.fill.fore_color.rgb = PRGB(0x05, 0x4F, 0x31)
    box3.line.color.rgb = PRGB(0x22, 0xC5, 0x5E); box3.line.width = Emu(12700)

    add_text(PI(9), col_top + PI(0.15), col_w - PI(0.4), PI(0.4),
             '📊 关键成果', 16, bold=True, color=(0x22,0xC5,0x5E))

    tb = slide.shapes.add_textbox(PI(9), col_top + PI(0.7), col_w - PI(0.4), col_h - PI(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    metrics = [
        ('2042', '条事件 (84→2042，+2330%)'),
        ('27 州', '60+ 城市，巴西全境覆盖'),
        ('30 min', '端到端实时延迟'),
        ('11 类', '犯罪类型分类'),
        ('$0', '/月运营成本'),
    ]
    for i, (n, v) in enumerate(metrics):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = n + '  '
        r.font.size = PP(15); r.font.bold = True
        r.font.color.rgb = PRGB(0x22,0xC5,0x5E); r.font.name = 'PingFang SC'
        r2 = p.add_run(); r2.text = v
        r2.font.size = PP(11); r2.font.color.rgb = PRGB(0xE5,0xE7,0xEB); r2.font.name = 'PingFang SC'
        p.space_after = PP(4)

    # === 底部：下一步 + 商业潜力 ===
    bottom_top = PI(4.55)
    bottom_h = PI(2.6)

    # 左：下一步
    box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PI(0.5), bottom_top, PI(6.1), bottom_h)
    box4.fill.solid(); box4.fill.fore_color.rgb = PRGB(0x1E, 0x3A, 0x8A)
    box4.line.color.rgb = PRGB(0x60, 0xA5, 0xFA); box4.line.width = Emu(12700)

    add_text(PI(0.7), bottom_top + PI(0.15), PI(5.7), PI(0.4),
             '🎯 下一步规划', 16, bold=True, color=(0x60,0xA5,0xFA))

    tb = slide.shapes.add_textbox(PI(0.7), bottom_top + PI(0.7), PI(5.7), bottom_h - PI(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    nexts = [
        ('🅰 高优', 'LLM 分类替代关键词 → 准确率 78% → 95%'),
        ('🅰 高优', '移动端 PWA + Push 推送通知'),
        ('🅱 中优', '接入警方 Twitter 官方账号 → 5 分钟级'),
        ('🅱 中优', '用户上报 UGC → 数据壁垒、社区粘性'),
        ('🅲 探索', 'AI 摘要 + 多语言（葡 → 英/中）'),
    ]
    for i, (lvl, t) in enumerate(nexts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = lvl + '  '
        r.font.size = PP(12); r.font.bold = True
        r.font.color.rgb = PRGB(0xFB,0xBF,0x24); r.font.name = 'PingFang SC'
        r2 = p.add_run(); r2.text = t
        r2.font.size = PP(12); r2.font.color.rgb = PRGB(0xE5,0xE7,0xEB); r2.font.name = 'PingFang SC'
        p.space_after = PP(4)

    # 右：商业潜力
    box5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PI(6.7), bottom_top, PI(6.1), bottom_h)
    box5.fill.solid(); box5.fill.fore_color.rgb = PRGB(0x58, 0x1C, 0x87)
    box5.line.color.rgb = PRGB(0xC0, 0x84, 0xFC); box5.line.width = Emu(12700)

    add_text(PI(6.9), bottom_top + PI(0.15), PI(5.7), PI(0.4),
             '💰 商业化机会', 16, bold=True, color=(0xC0,0x84,0xFC))

    tb = slide.shapes.add_textbox(PI(6.9), bottom_top + PI(0.7), PI(5.7), bottom_h - PI(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    biz = [
        ('广告', '本地商家按城市/邮编精准投放（餐饮、健身、地产）'),
        ('B2B', '保险公司车险定价 / 地产中介社区评估数据接口'),
        ('订阅', '"我家 1km 内的实时推送"高级用户付费'),
        ('数据', 'API 售卖给媒体研究机构、政府智库'),
    ]
    for i, (k, v) in enumerate(biz):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = k + '  '
        r.font.size = PP(12); r.font.bold = True
        r.font.color.rgb = PRGB(0xC0,0x84,0xFC); r.font.name = 'PingFang SC'
        r2 = p.add_run(); r2.text = v
        r2.font.size = PP(12); r2.font.color.rgb = PRGB(0xE5,0xE7,0xEB); r2.font.name = 'PingFang SC'
        p.space_after = PP(5)

    # 底部 footer
    add_text(PI(0.5), PI(7.15), PI(12.3), PI(0.3),
             '📍 演示：crime-map-recife.vercel.app   ·   📅 2026.06   ·   🚀 Crime Map Brasil MVP',
             10, color=(0x9C,0xA3,0xAF), align='center')

    out = OUT_DIR / '巴西犯罪地图_一页汇报.pptx'
    prs.save(out)
    print(f'✅ PPT:  {out}')


if __name__ == '__main__':
    make_docx()
    make_pptx()
    print('\n🎉 全部生成完毕！文件位于：', OUT_DIR)
